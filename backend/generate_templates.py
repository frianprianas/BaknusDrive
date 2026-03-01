import base64
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

# DOCX
doc = Document()
doc.save('empty.docx')

# XLSX
wb = Workbook()
wb.save('empty.xlsx')

# PPTX
prs = Presentation()
prs.slides.add_slide(prs.slide_layouts[0])
prs.save('empty.pptx')

# Encode to base64
def get_b64(filename):
    with open(filename, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")

docx_b64 = get_b64("empty.docx")
xlsx_b64 = get_b64("empty.xlsx")
pptx_b64 = get_b64("empty.pptx")

# Generate GO file
go_code = f"""package main

import (
\t"encoding/base64"
)

const BlankDocxBase64 = "{docx_b64}"
const BlankXlsxBase64 = "{xlsx_b64}"
const BlankPptxBase64 = "{pptx_b64}"

func CreateEmptyDocx() ([]byte, error) {{
\treturn base64.StdEncoding.DecodeString(BlankDocxBase64)
}}

func CreateEmptyXlsx() ([]byte, error) {{
\treturn base64.StdEncoding.DecodeString(BlankXlsxBase64)
}}

func CreateEmptyPptx() ([]byte, error) {{
\treturn base64.StdEncoding.DecodeString(BlankPptxBase64)
}}
"""

with open("office_templates.go", "w") as f:
    f.write(go_code)

print("Generated office_templates.go successfully.")
